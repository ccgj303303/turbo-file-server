"""
Turbo File Server — Railway v4
Genera PPTX y XLSX. Analiza CSV enviado como texto desde n8n.
Endpoints: POST /generate, POST /analyze, GET /health
"""

from flask import Flask, request, jsonify
import os, io, base64, traceback
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

API_TOKEN = os.environ.get("TURBO_API_TOKEN", "changeme123")

def check_auth(req):
    token = req.headers.get("X-Turbo-Token") or req.args.get("token")
    return token == API_TOKEN


# ══════════════════════════════════════════════════════════════════════════
# PPTX GENERATOR
# ══════════════════════════════════════════════════════════════════════════
def generate_pptx(payload):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Paleta Finvivir ──────────────────────────────────────────────────────
    C = {
        "blue":        "3B8DBD",
        "purple":      "7B3F7A",
        "dark_blue":   "1A3A5C",
        "gray":        "8A8A8A",
        "light_blue":  "EAF4FB",
        "light_gray":  "F5F5F7",
        "white":       "FFFFFF",
    }

    def rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def rect(slide, x, y, w, h, color):
        s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(color); s.line.fill.background()
        return s

    def textbox(slide, text, x, y, w, h, size,
                bold=False, italic=False, color="dark_blue",
                align=PP_ALIGN.LEFT, font="Calibri Light"):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; run = p.add_run()
        run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.italic = italic; run.font.name = font
        run.font.color.rgb = rgb(C[color])
        p.alignment = align
        return tb

    blank = prs.slide_layouts[6]

    # ── PORTADA ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(C["white"])

    # Panel derecho azul (45 % del ancho)
    rect(slide, 7.33, 0,    6.0,  7.5,  C["blue"])
    # Barra morada delgada de separación
    rect(slide, 7.13, 0,    0.2,  7.5,  C["purple"])
    # Acento morado inferior izquierdo
    rect(slide, 0,    6.9,  3.0,  0.6,  C["purple"])
    # Cuadrado decorativo blanco en panel derecho
    rect(slide, 11.8, 5.8,  1.2,  1.2,  C["dark_blue"])

    # Título
    textbox(slide, payload.get("title", "Presentación"),
            0.6, 1.6, 6.3, 2.2, 42, bold=True, color="dark_blue", font="Calibri")
    # Subtítulo
    textbox(slide, payload.get("subtitle", ""),
            0.6, 4.0, 6.2, 0.9, 18, color="gray", font="Calibri Light")
    # Slogan en panel derecho
    textbox(slide, "Abrazamos tu Futuro",
            7.7, 6.5, 5.3, 0.7, 13, italic=True, color="white",
            align=PP_ALIGN.RIGHT, font="Calibri Light")

    # ── SLIDES DE CONTENIDO ──────────────────────────────────────────────────
    for idx, slide_data in enumerate(payload.get("slides", [])):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(C["white"])

        # Barra azul superior
        rect(slide, 0,    0,    13.33, 0.07, C["blue"])
        # Franja morada izquierda
        rect(slide, 0,    0.07, 0.07,  7.08, C["purple"])
        # Zona de header gris muy suave
        rect(slide, 0.07, 0.07, 13.26, 1.05, C["light_gray"])

        # Título del slide
        textbox(slide, slide_data.get("title", ""),
                0.55, 0.14, 11.8, 0.9, 26, bold=True,
                color="dark_blue", font="Calibri")

        # Número de página
        textbox(slide, str(idx + 2),
                12.3, 0.18, 0.8, 0.5, 11, color="gray",
                align=PP_ALIGN.RIGHT, font="Calibri Light")

        # Footer
        rect(slide, 0, 7.15, 13.33, 0.35, C["light_gray"])
        rect(slide, 0, 7.13, 13.33, 0.02, C["blue"])
        textbox(slide, "Finvivir  ·  Abrazamos tu Futuro",
                0.55, 7.18, 9.0, 0.28, 9, color="gray", font="Calibri Light")

        # Bullets
        bullets = slide_data.get("bullets", [])
        if bullets:
            y = 1.38
            for bullet in bullets:
                is_sub = bullet.startswith("  ") or bullet.startswith("- ")
                text   = bullet.lstrip(" -").strip()

                if is_sub:
                    # Sub-bullet: línea morada delgada + texto gris
                    rect(slide, 0.75, y + 0.13, 0.18, 0.03, C["purple"])
                    textbox(slide, text, 1.1, y, 11.5, 0.44,
                            13, color="gray", font="Calibri Light")
                    y += 0.48
                else:
                    # Bullet principal: tarjeta con borde izquierdo azul
                    rect(slide, 0.45, y,       0.06, 0.52, C["blue"])
                    rect(slide, 0.51, y,       11.37, 0.52, C["light_gray"])
                    textbox(slide, text, 0.72, y + 0.02, 11.0, 0.48,
                            16, bold=True, color="dark_blue", font="Calibri")
                    y += 0.66

        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # ── SLIDE FINAL ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(C["blue"])

    # Panel morado izquierdo
    rect(slide, 0,    0,    4.8,  7.5,  C["purple"])
    # Cuadrado decorativo azul oscuro esquina superior derecha
    rect(slide, 11.5, 0,    1.83, 1.6,  C["dark_blue"])
    # Acento blanco pequeño
    rect(slide, 4.8,  6.5,  8.53, 0.06, C["white"])

    # Texto principal
    textbox(slide, "Abrazamos",
            5.2, 2.0, 7.7, 1.3, 52, bold=True,
            color="white", font="Calibri")
    textbox(slide, "tu Futuro",
            5.2, 3.2, 7.7, 1.3, 52,
            color="white", font="Calibri Light")

    # Título en panel morado
    textbox(slide, payload.get("title", ""),
            0.4, 5.6, 4.0, 1.2, 13, italic=True,
            color="white", font="Calibri Light")

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════════════════
# XLSX GENERATOR
# ══════════════════════════════════════════════════════════════════════════
def generate_xlsx(payload):
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    HEADER_FILL = PatternFill("solid", fgColor="3B8DBD")   # Azul Finvivir
    ACCENT_FILL = PatternFill("solid", fgColor="7B3F7A")   # Morado Finvivir
    ALT_FILL    = PatternFill("solid", fgColor="EAF4FB")   # Azul muy claro
    TOTAL_FILL  = PatternFill("solid", fgColor="F3EDF7")   # Lila muy claro
    HEADER_FONT = Font(name="Arial", bold=True, color="FFFFFF", size=11)
    TITLE_FONT  = Font(name="Arial", bold=True, color="3B8DBD", size=14)
    DATA_FONT   = Font(name="Arial", size=10)
    TOTAL_FONT  = Font(name="Arial", bold=True, color="3B8DBD", size=11)
    CENTER = Alignment(horizontal="center", vertical="center")
    LEFT   = Alignment(horizontal="left",   vertical="center")
    thin   = Side(style="thin", color="D0D7DE")
    BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)
    for sheet_def in payload.get("sheets", []):
        ws = wb.create_sheet(title=sheet_def.get("name", "Report")[:31])
        headers  = sheet_def.get("headers", [])
        rows     = sheet_def.get("rows", [])
        do_total = sheet_def.get("totals", False)
        widths   = sheet_def.get("col_widths", [])
        ws.row_dimensions[1].height = 30
        title_cell = ws.cell(row=1, column=1, value=payload.get("title", "Report"))
        title_cell.font = TITLE_FONT; title_cell.alignment = LEFT
        if headers:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
        for col in range(1, len(headers)+1):
            ws.cell(row=2, column=col, value="").fill = ACCENT_FILL
        ws.row_dimensions[3].height = 22
        for col_idx, header in enumerate(headers, 1):
            c = ws.cell(row=3, column=col_idx, value=header)
            c.font = HEADER_FONT; c.fill = HEADER_FILL; c.alignment = CENTER; c.border = BORDER
        for row_idx, row_data in enumerate(rows, 4):
            ws.row_dimensions[row_idx].height = 18
            for col_idx, value in enumerate(row_data, 1):
                c = ws.cell(row=row_idx, column=col_idx, value=value)
                c.font = DATA_FONT; c.border = BORDER
                c.alignment = CENTER if isinstance(value, (int,float)) else LEFT
                if row_idx % 2 == 0: c.fill = ALT_FILL
        if do_total and rows:
            total_row = len(rows) + 4
            ws.row_dimensions[total_row].height = 22
            ws.cell(row=total_row, column=1, value="TOTAL").font = TOTAL_FONT
            ws.cell(row=total_row, column=1).fill = TOTAL_FILL
            for col_idx in range(2, len(headers)+1):
                col_letter = get_column_letter(col_idx)
                c = ws.cell(row=total_row, column=col_idx,
                            value=f"=SUM({col_letter}4:{col_letter}{total_row-1})")
                c.font = TOTAL_FONT; c.fill = TOTAL_FILL; c.alignment = CENTER; c.border = BORDER
        for col_idx, header in enumerate(headers, 1):
            col_letter = get_column_letter(col_idx)
            ws.column_dimensions[col_letter].width = (
                widths[col_idx-1] if widths and col_idx <= len(widths)
                else max(len(str(header))+4, 14)
            )
        ws.freeze_panes = "A4"
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════════════════════════
# SHARED: build_analysis
# ══════════════════════════════════════════════════════════════════════════
def build_analysis(df, filename, instruction):
    df.columns = [
        str(c) if not str(c).startswith("Unnamed") else f"Col_{i+1}"
        for i, c in enumerate(df.columns)
    ]
    df = df.dropna(how="all")
    rows, columns = df.shape
    insights = []
    numeric_cols = df.select_dtypes(include="number").columns.tolist()
    for col in numeric_cols[:5]:
        insights.append(
            f"{col}: Total={df[col].sum():,.2f} | "
            f"Avg={df[col].mean():,.2f} | "
            f"Max={df[col].max():,.2f} | "
            f"Min={df[col].min():,.2f}"
        )
    nulls = df.isnull().sum()
    null_cols = nulls[nulls > 0]
    if len(null_cols) > 0:
        insights.append(f"⚠️ Valores nulos en: {', '.join(null_cols.index.tolist()[:5])}")
    else:
        insights.append("✅ Sin valores nulos")
    insights.append(f"Total de registros: {rows:,}")
    data_preview = df.head(5).to_string(index=False, max_cols=8)
    col_list = ", ".join(df.columns.tolist()[:10])
    summary = f"Columnas: {col_list}" + (
        f" ... (+{len(df.columns)-10} más)" if len(df.columns) > 10 else ""
    )
    return {
        "success":      True,
        "filename":     filename,
        "rows":         rows,
        "columns":      columns,
        "col_names":    df.columns.tolist(),
        "insights":     insights,
        "summary":      summary,
        "data_preview": data_preview,
        "instruction":  instruction
    }


# ══════════════════════════════════════════════════════════════════════════
# ROUTES
# ══════════════════════════════════════════════════════════════════════════
@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "service": "Turbo File Server v4"})


@app.route("/generate", methods=["POST"])
def generate():
    if not check_auth(request):
        return jsonify({"error": "Unauthorized"}), 401
    try:
        body      = request.get_json(force=True)
        file_type = body.get("type", "").lower()
        payload   = body.get("payload", {})
        filename  = body.get("filename", "output")
        if file_type == "pptx":
            file_bytes = generate_pptx(payload)
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ext  = "pptx"
        elif file_type == "xlsx":
            file_bytes = generate_xlsx(payload)
            mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ext  = "xlsx"
        else:
            return jsonify({"error": f"Unknown type: {file_type}"}), 400
        encoded = base64.b64encode(file_bytes).decode("utf-8")
        return jsonify({
            "success": True, "type": ext, "filename": f"{filename}.{ext}",
            "mime": mime, "size_bytes": len(file_bytes), "data_base64": encoded
        })
    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route("/analyze", methods=["POST"])
def analyze():
    if not check_auth(request):
        return jsonify({"error": "Unauthorized"}), 401
    try:
        body = request.get_json(force=True, silent=True)
        if not body:
            return jsonify({"error": "No se recibió payload JSON"}), 400

        # ── Recibir CSV como texto desde n8n ─────────────────────────────
        csv_text    = body.get("csv_text", "")
        instruction = body.get("instruction", "Analiza este archivo")
        filename    = body.get("filename", "archivo.csv")

        if not csv_text:
            return jsonify({"error": "No se recibió csv_text"}), 400

        # ── Leer CSV desde el string de texto ────────────────────────────
        df = None
        errors = []

        for enc in ["utf-8", "latin-1", "cp1252"]:
            try:
                buf = io.StringIO(csv_text)
                df = pd.read_csv(buf)
                if len(df.columns) > 0 and len(df) > 0:
                    break
            except Exception as e:
                errors.append(f"csv/{enc}: {e}")

        if df is None or len(df.columns) == 0:
            return jsonify({
                "error": "No se pudo leer el CSV.",
                "attempts": errors[:3]
            }), 500

        return jsonify(build_analysis(df, filename, instruction))

    except Exception as e:
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port, debug=False)
